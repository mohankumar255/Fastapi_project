from fastapi import FastAPI, File, UploadFile, HTTPException, Depends
from fastapi.security import HTTPBasic, HTTPBasicCredentials
from pydantic import BaseModel
from pymongo import MongoClient
from uuid import uuid4
import os
from typing import List
import docx
import pptx
from PyPDF2 import PdfReader
import requests

app = FastAPI()
security = HTTPBasic()

# MongoDB configuration
client = MongoClient("mongodb://localhost:27017/")
db = client.file_database
collection = db.files

# Storage folder
STORAGE_FOLDER = "./storage"
os.makedirs(STORAGE_FOLDER, exist_ok=True)


# Authentication
def authenticate(credentials):
    correct_username = "user"
    correct_password = "password"
    if credentials['username'] != correct_username or credentials['password'] != correct_password:
        raise HTTPException(status_code=401, detail="Invalid credentials")
    return credentials['username']


# Models
class FileInfo(BaseModel):
    file_id: str
    file_name: str
    file_summary: str


# Helper function to extract text from files
def extract_text(file: UploadFile):
    if file.filename.endswith(".docx"):
        doc = docx.Document(file.file)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    elif file.filename.endswith(".pptx"):
        presentation = pptx.Presentation(file.file)
        text = "\n".join(
            [shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif file.filename.endswith(".pdf"):
        pdf = PdfReader(file.file)
        text = "\n".join([page.extract_text() for page in pdf.pages])
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")
    return text


# Helper function to summarize text using Predibase
def summarize_text(text: str):
    response = requests.post(
        "https://api.predibase.com/summarize",
        data={"text": text, "lines": 3},
        headers={"Authorization": "Bearer pb_duf3uTeDU-mSCuYER7OP7Ql"}
    )
    if response.status_code == 200:
        return response.json().get("summary")
    else:
        raise HTTPException(status_code=500, detail="Error summarizing text")


# Endpoints
@app.post("/v1/files", response_model=FileInfo)
def upload_file(file: UploadFile,username:str,password:str):
    credentials = {'username': username, 'password': password}
    authentication = authenticate(credentials)
    if str(authentication) == username:
        if not file.filename.endswith((".docx", ".pptx", ".pdf")):
            raise HTTPException(status_code=400, detail="Invalid file type")
        existing_file = collection.find_one({"file_name": file.filename})
        if existing_file:
            raise HTTPException(status_code=400, detail="File already uploaded")
        file_id = str(uuid4())
        file_path = os.path.join(STORAGE_FOLDER, file_id)
        with open(file_path, "wb") as f:
            f.write(file.read())
        text = extract_text(file)
        summary = summarize_text(text)
        file_info = {
            "file_id": file_id,
            "file_name": file.filename,
            "file_summary": summary
        }
        collection.insert_one(file_info)
        return file_info
    return authentication

@app.get("/v1/files", response_model=List[str])
def list_files(username:str,password:str):
    credentials={'username':username,'password':password}
    authentication = authenticate(credentials)
    if str(authentication)==username:
        files = collection.find({}, {"_id": 0, "file_id": 1})
        return [file["file_id"] for file in files]
    return authentication

@app.get("/v1/files/{file_id}", response_model=FileInfo)
def get_file_summary(file_id: str, username:str,password:str):
    credentials = {'username': username, 'password': password}
    authentication = authenticate(credentials)
    if str(authentication)!=username:
        file_info = collection.find_one({"file_id": file_id}, {"_id": 0})
        if not file_info:
            raise HTTPException(status_code=404, detail="File not found")
        return file_info
    return authentication


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app,host = '127.0.0.1',port=8000)
